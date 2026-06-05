#!/usr/bin/env python3
"""
Analyze a PowerPoint presentation and provide improvement suggestions.
Usage: python analyze_pptx.py <file.pptx> [--verbose]
Checks: text density, consistency, images, readability, structure.
"""
import json, sys, argparse
from pathlib import Path
from collections import Counter


def analyze(filepath, verbose=False):
    from pptx import Presentation
    from pptx.util import Pt, Emu

    prs = Presentation(filepath)
    
    issues = []
    stats = {
        "total_slides": len(prs.slides),
        "slides_with_images": 0,
        "slides_with_charts": 0,
        "slides_with_tables": 0,
        "slides_with_notes": 0,
        "total_shapes": 0,
        "total_text_chars": 0,
        "fonts_used": Counter(),
        "font_sizes_used": Counter(),
        "layouts_used": Counter(),
        "avg_text_per_slide": 0,
        "empty_slides": 0,
        "text_heavy_slides": [],
        "text_light_slides": [],
    }
    
    slide_texts = []
    
    for slide_number, slide in enumerate(prs.slides, start=1):
        stats["layouts_used"][slide.slide_layout.name] += 1
        slide_text_len = 0
        slide_has_image = False
        slide_has_chart = False
        slide_has_table = False
        has_title = False
        small_font_warnings = set()
        
        for shape in slide.shapes:
            stats["total_shapes"] += 1
            
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                slide_text_len += len(text)
                
                try:
                    if shape.placeholder_format.type in (1, 2, 15, 16):
                        has_title = bool(text)
                except Exception:
                    pass
                
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            stats["fonts_used"][run.font.name] += 1
                        if run.font.size:
                            size_pt = round(run.font.size / 12700, 0)
                            stats["font_sizes_used"][f"{size_pt}pt"] += 1
                            top_inches = shape.top / 914400 if shape.top else 0
                            height_inches = shape.height / 914400 if shape.height else 0
                            is_footer_like = top_inches >= 6.65 and height_inches <= 0.4
                            warning_key = (shape.name, size_pt)
                            if size_pt < 12 and len(run.text.strip()) > 20 and not is_footer_like and warning_key not in small_font_warnings:
                                small_font_warnings.add(warning_key)
                                issues.append({
                                    "slide": slide_number,
                                    "severity": "warning",
                                    "issue": f"Small font ({size_pt}pt) on slide {slide_number} - may be hard to read in presentation",
                                    "shape": shape.name,
                                })
            
            try:
                _ = shape.image.content_type
                slide_has_image = True
            except Exception:
                pass
            
            if shape.has_chart:
                slide_has_chart = True
            
            if shape.has_table:
                slide_has_table = True
        
        stats["total_text_chars"] += slide_text_len
        slide_texts.append(slide_text_len)
        
        if slide_has_image:
            stats["slides_with_images"] += 1
        if slide_has_chart:
            stats["slides_with_charts"] += 1
        if slide_has_table:
            stats["slides_with_tables"] += 1
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            stats["slides_with_notes"] += 1
        
        if slide_text_len == 0 and not slide_has_image and not slide_has_chart:
            stats["empty_slides"] += 1
            issues.append({"slide": slide_number, "severity": "info", "issue": f"Slide {slide_number} appears empty"})
        
        if slide_text_len > 500:
            stats["text_heavy_slides"].append(slide_number)
            issues.append({
                "slide": slide_number, "severity": "warning",
                "issue": f"Slide {slide_number} has {slide_text_len} chars - consider splitting or reducing text"
            })
        
        if slide_text_len > 0 and slide_text_len < 30 and not slide_has_image and not slide_has_chart:
            stats["text_light_slides"].append(slide_number)
        
        if not has_title and slide_number > 1:
            issues.append({"slide": slide_number, "severity": "info", "issue": f"Slide {slide_number} has no title"})
    
    if slide_texts:
        stats["avg_text_per_slide"] = round(sum(slide_texts) / len(slide_texts), 0)
    
    # Font consistency check
    if len(stats["fonts_used"]) > 3:
        issues.append({
            "slide": "all", "severity": "warning",
            "issue": f"Using {len(stats['fonts_used'])} different fonts ({', '.join(stats['fonts_used'].keys())}) - consider limiting to 2-3 for consistency"
        })
    
    # Notes coverage
    notes_pct = round(stats["slides_with_notes"] / stats["total_slides"] * 100) if stats["total_slides"] > 0 else 0
    if notes_pct < 50 and stats["total_slides"] > 3:
        issues.append({
            "slide": "all", "severity": "info",
            "issue": f"Only {notes_pct}% of slides have speaker notes - consider adding notes for presentation delivery"
        })
    
    # Visual balance
    img_pct = round(stats["slides_with_images"] / stats["total_slides"] * 100) if stats["total_slides"] > 0 else 0
    if img_pct < 20 and stats["total_slides"] > 5:
        issues.append({
            "slide": "all", "severity": "suggestion",
            "issue": "Less than 20% of slides have images - consider adding visuals for engagement"
        })
    
    # Convert counters to dicts for JSON
    stats["fonts_used"] = dict(stats["fonts_used"])
    stats["font_sizes_used"] = dict(stats["font_sizes_used"])
    stats["layouts_used"] = dict(stats["layouts_used"])
    
    return {"stats": stats, "issues": issues}


def main():
    parser = argparse.ArgumentParser(description="Analyze PowerPoint presentation")
    parser.add_argument("file", help="PowerPoint file path")
    parser.add_argument("--verbose", action="store_true", help="Include all details")
    args = parser.parse_args()
    
    if not Path(args.file).exists():
        print(json.dumps({"error": f"File not found: {args.file}"}))
        sys.exit(1)
    
    result = analyze(args.file, args.verbose)
    print(json.dumps(result, indent=2, default=str))


if __name__ == "__main__":
    main()
