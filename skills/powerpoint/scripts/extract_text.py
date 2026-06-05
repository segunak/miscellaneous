#!/usr/bin/env python3
"""
Extract all text from a PowerPoint file as structured JSON or plain text.
Usage: python extract_text.py <file.pptx> [--format json|text|markdown] [--notes]
"""
import json, sys, argparse
from pathlib import Path

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8")


def extract(filepath, fmt="markdown", include_notes=False):
    from pptx import Presentation

    prs = Presentation(filepath)
    slides_data = []
    
    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_content = {"slide": slide_number, "title": "", "body": [], "notes": ""}
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                # Detect title
                try:
                    ph_type = shape.placeholder_format.type
                    if ph_type in (1, 2, 15, 16):  # TITLE, CENTER_TITLE, SUBTITLE types
                        slide_content["title"] = text
                        continue
                except Exception:
                    pass
                slide_content["body"].append(text)
            
            if shape.has_table:
                table_text = []
                for row in shape.table.rows:
                    table_text.append([cell.text for cell in row.cells])
                slide_content["body"].append({"table": table_text})
        
        if include_notes and slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_content["notes"] = notes_text
        
        slides_data.append(slide_content)
    
    if fmt == "json":
        return json.dumps(slides_data, indent=2, default=str)
    elif fmt == "text":
        lines = []
        for s in slides_data:
            lines.append(f"=== Slide {s['slide']} ===")
            if s["title"]:
                lines.append(f"Title: {s['title']}")
            for b in s["body"]:
                if isinstance(b, dict) and "table" in b:
                    for row in b["table"]:
                        lines.append(" | ".join(row))
                else:
                    lines.append(b)
            if s.get("notes"):
                lines.append(f"[Notes: {s['notes']}]")
            lines.append("")
        return "\n".join(lines)
    else:  # markdown
        lines = []
        for s in slides_data:
            lines.append(f"## Slide {s['slide']}")
            if s["title"]:
                lines.append(f"### {s['title']}")
            for b in s["body"]:
                if isinstance(b, dict) and "table" in b:
                    headers = b["table"][0]
                    lines.append("| " + " | ".join(headers) + " |")
                    lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                    for row in b["table"][1:]:
                        lines.append("| " + " | ".join(row) + " |")
                else:
                    lines.append(b)
            if s.get("notes"):
                lines.append(f"\n> **Notes:** {s['notes']}")
            lines.append("")
        return "\n".join(lines)


def main():
    parser = argparse.ArgumentParser(description="Extract text from PowerPoint")
    parser.add_argument("file", help="PowerPoint file path")
    parser.add_argument("--format", choices=["json", "text", "markdown"], default="markdown")
    parser.add_argument("--notes", action="store_true", help="Include speaker notes")
    args = parser.parse_args()
    
    if not Path(args.file).exists():
        print(f"Error: File not found: {args.file}")
        sys.exit(1)
    
    print(extract(args.file, args.format, args.notes))


if __name__ == "__main__":
    main()
